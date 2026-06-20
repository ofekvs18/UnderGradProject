"""Build the academic poster for the biomarker-pipeline project.

Assets used:
  results/poster/svg_framework_diagram.png   — SVG rendered by PyMuPDF (run once)
  results/poster/svg_feature_consensus.png   — SVG rendered by PyMuPDF
  results/ra_aucpr_lift.png                  — real Matplotlib chart with CI + complexity

Run from the repo root:
    ../.venv/Scripts/python.exe results/poster/build_poster.py
"""

import io
import os
import fitz                      # PyMuPDF — SVG → PNG
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Colors  (match the SVG palette)
# ---------------------------------------------------------------------------
TEAL       = RGBColor(0x1A, 0x7A, 0x8A)
TEAL_LIGHT = RGBColor(0xD6, 0xEA, 0xF0)
ORANGE     = RGBColor(0xE8, 0x77, 0x3D)
ORANGE_LT  = RGBColor(0xFD, 0xEB, 0xD7)
NAVY       = RGBColor(0x2C, 0x3E, 0x50)
GRAY       = RGBColor(0x7F, 0x8C, 0x8D)
GRAY_LT    = RGBColor(0xF0, 0xF3, 0xF4)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------------------
# SVG → PNG conversion (skip if PNG already fresh)
# ---------------------------------------------------------------------------
POSTER_DIR = "results/poster"
SVG_FILES  = ["svg_framework_diagram", "svg_feature_consensus", "svg_ra_aucpr_lift"]

def ensure_pngs(scale: int = 4):
    for name in SVG_FILES:
        svg = os.path.join(POSTER_DIR, name + ".svg")
        png = os.path.join(POSTER_DIR, name + ".png")
        if not os.path.exists(png) or os.path.getmtime(svg) > os.path.getmtime(png):
            doc = fitz.open(svg)
            pix = doc[0].get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
            with open(png, "wb") as f:
                f.write(pix.tobytes("png"))
            doc.close()
            print(f"Rendered {name}.png")

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def I(v):
    """Inches → EMU."""
    return int(v * 914400)

def _set_rtl(p):
    p._p.get_or_add_pPr().set("rtl", "1")

def add_rect(slide, x, y, w, h, fill, line=None, line_pt=None):
    s = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if line_pt:
            s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def add_textbox(slide, x, y, w, h, text="", size=22, bold=False,
                color=NAVY, align=PP_ALIGN.RIGHT, rtl=True, italic=False):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if rtl:
        _set_rtl(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def section_header(slide, x, y, w, h, title, color, size=26):
    """Solid-fill bar with centred white title."""
    bar = add_rect(slide, x, y, w, h, color)
    tf  = bar.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = WHITE

def card(slide, x, y, w, h, icon, title, body,
         bg=GRAY_LT, accent=TEAL, title_size=24, body_size=20):
    add_rect(slide, x, y, w, h, bg, accent, 1.0)
    add_textbox(slide, x + w - 0.75, y + 0.10, 0.65, 0.55,
                icon, 32, align=PP_ALIGN.CENTER, rtl=False)
    add_textbox(slide, x + 0.15, y + 0.12, w - 0.95, 0.52,
                title, title_size, bold=True, color=accent, align=PP_ALIGN.RIGHT)
    add_textbox(slide, x + 0.15, y + 0.70, w - 0.30, h - 0.82,
                body, body_size, color=NAVY, align=PP_ALIGN.RIGHT)

def add_image(slide, path, x, y, w):
    """Add picture at (x,y) with given width; height scales to preserve aspect."""
    with Image.open(path) as im:
        iw, ih = im.size
    h = w * ih / iw
    slide.shapes.add_picture(path, I(x), I(y), I(w), I(h))
    return h   # return actual height in inches

def labeled_row(slide, x, y, w, h, label, text, label_w, fill):
    add_rect(slide, x, y, w, h, fill)
    add_textbox(slide, x + 0.12, y + 0.10, label_w, h - 0.15,
                label, 21, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    add_textbox(slide, x + label_w + 0.20, y + 0.10, w - label_w - 0.32, h - 0.15,
                text, 20, color=WHITE, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# Main build
# ---------------------------------------------------------------------------
def build():
    ensure_pngs()

    prs   = Presentation("results/poster/פוסטר.pptx")
    slide = prs.slides[0]
    W     = prs.slide_width.inches    # 35.43
    H     = prs.slide_height.inches   # 47.25

    MARGIN    = 0.55
    GAP       = 0.40
    col_w     = (W - 2 * MARGIN - 2 * GAP) / 3   # ≈ 11.18"
    col_x     = [MARGIN + i * (col_w + GAP) for i in range(3)]
    Y0        = 7.35     # below BGU header
    FOOTER_H  = 1.05
    CONTENT_H = H - Y0 - FOOTER_H   # ≈ 38.85"
    SH        = 0.55     # section-header bar height
    PAD       = 0.25     # inner padding

    # ── Image dimensions (aspect-ratio driven) ──────────────────────────────
    IMG_W     = col_w - 0.40   # image width = column minus side padding
    fw_h      = IMG_W * 420 / 900    # framework diagram  900×420 viewBox → ≈ 5.19"
    cons_h    = IMG_W * 320 / 620    # feature consensus  620×320 viewBox → ≈ 5.75"
    ra_h      = IMG_W * 490 / 760    # RA lift SVG        760×490 viewBox → ≈ 7.22"

    # ── Footer ──────────────────────────────────────────────────────────────
    add_rect(slide, 0, H - FOOTER_H, W, FOOTER_H, NAVY)
    add_textbox(slide, 0.5, H - FOOTER_H + 0.12, W - 1.0, 0.75,
                "Ofek Vilozny  |  מנחה: פרופ' נעה דגן  |  הנדסת מערכות מידע, BGU  |  2026",
                size=20, color=TEAL_LIGHT, align=PP_ALIGN.CENTER, rtl=False)

    # ══════════════════════════════════════════════════════════════════════
    # COLUMN 1 — Motivation · Problem · Metrics
    # ══════════════════════════════════════════════════════════════════════
    cx = col_x[0]

    # ── Section 1: מוטיבציה ורקע ─────────────────────────────────────────
    S1_Y, S1_H = Y0, 12.8
    add_rect(slide, cx, S1_Y, col_w, S1_H, GRAY_LT, ORANGE, 1.5)
    section_header(slide, cx, S1_Y, col_w, SH, "מוטיבציה ורקע", ORANGE)

    for i, (icon, title, body, bg, accent) in enumerate([
        ("⚠", "האתגר",
         "אבחון מחלות אוטואימוניות (זאבת, קרוהן, RA, סוכרת, פסוריאזיס) "
         "מסתמך על בדיקות יקרות, מורכבות ומאוחרות.",
         ORANGE_LT, ORANGE),
        ("💡", "ההזדמנות",
         "ספירת דם מלאה (CBC) היא בדיקה זולה, שגרתית וזמינה בכל מפגש "
         "רפואי — אך הפוטנציאל שלה כביומרקר אבחנתי טרם מוצה.",
         TEAL_LIGHT, TEAL),
        ("🔬", "החזון",
         "מחקר זה בוחן האם ניתן לגזור נוסחאות מתמטיות פשוטות מערכי CBC "
         "לצורך חיזוי ואבחון מוקדם של 6 מחלות.",
         GRAY_LT, NAVY),
    ]):
        card(slide, cx + 0.20, S1_Y + SH + 0.20 + i * 3.90, col_w - 0.40, 3.70,
             icon, title, body, bg=bg, accent=accent)

    # ── Section 2: הגדרת הבעיה ───────────────────────────────────────────
    S2_Y = S1_Y + S1_H + GAP
    S2_H = 10.5
    add_rect(slide, cx, S2_Y, col_w, S2_H, GRAY_LT, TEAL, 1.5)
    section_header(slide, cx, S2_Y, col_w, SH, "הגדרת הבעיה", TEAL)

    # Research-question banner
    q_y = S2_Y + SH + PAD
    add_rect(slide, cx + 0.20, q_y, col_w - 0.40, 1.75, NAVY)
    add_textbox(slide, cx + 0.28, q_y + 0.12, col_w - 0.56, 1.52,
                "שאלת המחקר: איזו שיטה לייצור נוסחאות מניבה את המודל המיטבי "
                "— וכיצד להשוות ביניהן בצורה הוגנת ומבוקרת?",
                size=21, color=WHITE, align=PP_ALIGN.RIGHT)

    ROW_H, LBL_W = 1.55, 1.65
    for i, (label, text, fill) in enumerate([
        ("קלט",     "14 מאפיינים מ-CBC: ספירת תאים + אחוזי ד\"ל", TEAL),
        ("יעד",     "נוסחה פשוטה מ-≤5 מאפיינים, שקולה ל-LR עם כל הנתונים", ORANGE),
        ("מחלות",   "RA · קרוהן · T1D · T2D · פסוריאזיס · זאבת (6)", NAVY),
        ("נתונים",  "MIMIC-IV v3.1  |  BigQuery  |  split 80/20 קבוע", GRAY),
    ]):
        labeled_row(slide, cx + 0.20, q_y + 1.95 + i * (ROW_H + 0.10),
                    col_w - 0.40, ROW_H, label, text, LBL_W, fill)

    # ── Section 3: נתונים ומטריקות ──────────────────────────────────────
    S3_Y = S2_Y + S2_H + GAP
    S3_H = CONTENT_H - (S3_Y - Y0)
    add_rect(slide, cx, S3_Y, col_w, S3_H, GRAY_LT, NAVY, 1.5)
    section_header(slide, cx, S3_Y, col_w, SH, "נתונים ומטריקות", NAVY)

    metric_h = (S3_H - SH - 0.30) / 4 - 0.12
    for i, (label, name, detail, fill) in enumerate([
        ("מטריקה ראשית",    "AUC-PR Lift",       "AUC-PR / prevalence  (class imbalance ~1%)", ORANGE),
        ("מטריקה משנית",    "AUC-ROC",           "רגישות + ספציפיות",                        TEAL),
        ("רווחי סמך",       "Bootstrap CI",      "n=1,000 iterations  |  95% CI",             NAVY),
        ("ולידציה חיצונית", "NHANES + EHRSHOT",  "Stanford OMOP CDM  (generalisation test)",  GRAY),
    ]):
        my = S3_Y + SH + 0.18 + i * (metric_h + 0.12)
        add_rect(slide, cx + 0.20, my, col_w - 0.40, metric_h, fill)
        add_textbox(slide, cx + 0.28, my + 0.08, col_w - 0.56, 0.44,
                    label, 17, color=WHITE, align=PP_ALIGN.RIGHT)
        add_textbox(slide, cx + 0.28, my + 0.48, col_w - 0.56, 0.55,
                    name, 23, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        add_textbox(slide, cx + 0.28, my + 1.00, col_w - 0.56, metric_h - 1.08,
                    detail, 18, color=TEAL_LIGHT, align=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════════════════
    # COLUMN 2 — Framework diagram SVG · Feature Consensus SVG
    # ══════════════════════════════════════════════════════════════════════
    cx = col_x[1]
    img_x = cx + 0.20   # image left edge (0.2" inner padding)

    # ── Section: ארכיטקטורת הפריימוורק ──────────────────────────────────
    F_H = CONTENT_H * 0.575    # ≈ 22.3"
    F_Y = Y0
    add_rect(slide, cx, F_Y, col_w, F_H, GRAY_LT, TEAL, 1.5)
    section_header(slide, cx, F_Y, col_w, SH, "ארכיטקטורת הפריימוורק", TEAL)

    # Framework diagram image
    fw_img_y = F_Y + SH + PAD
    add_image(slide, os.path.join(POSTER_DIR, "svg_framework_diagram.png"),
              img_x, fw_img_y, IMG_W)

    # 4 method pills in 2×2 grid below the image
    pill_top  = fw_img_y + fw_h + 0.45
    pill_w    = (IMG_W - 0.20) / 2
    pill_h    = (F_H - SH - PAD - fw_h - 0.55) / 2 - 0.15

    for i, (label, name, detail, fill) in enumerate([
        ("M1", "Threshold Opt.",   "Youden's Index\nLiterature cutoffs\n1 feature",       NAVY),
        ("M2", "Random Search",    "10,000 formulas\n±×÷ √ log operators\nExhaustive",    GRAY),
        ("M3", "Genetic Prog.  🏆","gplearn  |  pop=500\ngen=100  |  parsimony=0.0001\nBest performer", ORANGE),
        ("M4", "LLM (Med-Gemma)", "Med-Gemma 4B\nSeeded prompts\nGPU cluster",            TEAL),
    ]):
        col_offset = (i % 2) * (pill_w + 0.20)
        row_offset = (i // 2) * (pill_h + 0.15)
        px = img_x + col_offset
        py = pill_top + row_offset
        add_rect(slide, px, py, pill_w, pill_h, fill)
        add_textbox(slide, px + 0.10, py + 0.10, 0.55, 0.48,
                    label, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
        add_textbox(slide, px + 0.70, py + 0.10, pill_w - 0.80, 0.48,
                    name, 19, bold=True, color=WHITE, align=PP_ALIGN.LEFT, rtl=False)
        add_textbox(slide, px + 0.10, py + 0.58, pill_w - 0.20, pill_h - 0.65,
                    detail, 17, color=TEAL_LIGHT, align=PP_ALIGN.LEFT, rtl=False)

    # ── Section: Cross-Method Feature Consensus ──────────────────────────
    FC_Y = F_Y + F_H + GAP
    FC_H = CONTENT_H - (FC_Y - Y0)
    add_rect(slide, cx, FC_Y, col_w, FC_H, GRAY_LT, ORANGE, 1.5)
    section_header(slide, cx, FC_Y, col_w, SH,
                   "Cross-Method Feature Consensus — RA", ORANGE)

    # Consensus heatmap image
    cons_img_y = FC_Y + SH + PAD
    add_image(slide, os.path.join(POSTER_DIR, "svg_feature_consensus.png"),
              img_x, cons_img_y, IMG_W)

    # Interpretation text below the heatmap
    interp_y = cons_img_y + cons_h + 0.30
    add_rect(slide, cx + 0.20, interp_y, col_w - 0.40, 0.55, TEAL)
    add_textbox(slide, cx + 0.28, interp_y + 0.08, col_w - 0.56, 0.42,
                "RDW · PLT · MCHC — consensus biomarkers across all 4 methods",
                size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)

    interp_items = [
        ("RDW",  "הכי עקבי: 85% ב-GP, 100% ב-Threshold, 64% ב-LLM"),
        ("PLT",  "מופיע ב-3/4 שיטות; GP 90%, Random 60%"),
        ("MCHC", "אינדיקטור לאנמיה — הוסק באופן עצמאי בכל שיטה"),
    ]
    for ii, (feat, desc) in enumerate(interp_items):
        iy = interp_y + 0.65 + ii * ((FC_H - SH - PAD - cons_h - 0.65 - 0.30) / 3 - 0.10)
        ih = (FC_H - SH - PAD - cons_h - 0.65 - 0.30) / 3 - 0.15
        feat_colors = [ORANGE, TEAL, NAVY]
        add_rect(slide, cx + 0.20, iy, 1.50, ih, feat_colors[ii])
        add_textbox(slide, cx + 0.20, iy, 1.50, ih,
                    feat, 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
        add_textbox(slide, cx + 1.78, iy + 0.05, col_w - 1.98, ih - 0.10,
                    desc, 19, color=NAVY, align=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════════════════
    # COLUMN 3 — RA results chart · Conclusions · External validation
    # ══════════════════════════════════════════════════════════════════════
    cx = col_x[2]
    img_x3 = cx + 0.20

    # ── Section: תוצאות — Rheumatoid Arthritis ───────────────────────────
    R_H = CONTENT_H * 0.445    # ≈ 17.3"
    R_Y = Y0
    add_rect(slide, cx, R_Y, col_w, R_H, GRAY_LT, ORANGE, 1.5)
    section_header(slide, cx, R_Y, col_w, SH,
                   "תוצאות — Rheumatoid Arthritis", ORANGE)

    # RA AUC-PR Lift SVG (real data, CI error bars, formula complexity)
    chart_y = R_Y + SH + PAD
    add_image(slide, os.path.join(POSTER_DIR, "svg_ra_aucpr_lift.png"), img_x3, chart_y, IMG_W)

    # Key findings below the chart
    findings_y = chart_y + ra_h + 0.35
    add_textbox(slide, cx + 0.15, findings_y, col_w - 0.30, 0.50,
                "ממצאים עיקריים", 23, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

    finding_h = (R_H - SH - PAD - ra_h - 0.45 - 0.55) / 4 - 0.10
    for fi, (text, bar_col) in enumerate([
        ("GP מנצח: 1.59× lift עם 4 features — עולה על LR (1.49×) שמשתמש ב-14",  ORANGE),
        ("LLM הכי יעיל: 3 features, 2 ops, lift 1.45× — כמעט שקול ל-LR",         TEAL),
        ("Random לא מספיק: 7 features, 18 ops → 1.32× (מתחת baseline)",           GRAY),
        ("Threshold נקודת ייחוס: feature יחיד (rbc<3.99) → 1.38×",               NAVY),
    ]):
        fy = findings_y + 0.55 + fi * (finding_h + 0.10)
        add_rect(slide, cx + 0.15, fy, 0.10, finding_h, bar_col)
        add_textbox(slide, cx + 0.35, fy + 0.05, col_w - 0.50, finding_h - 0.08,
                    text, 19, color=NAVY, align=PP_ALIGN.RIGHT)

    # ── Section: מסקנות ──────────────────────────────────────────────────
    C_Y = R_Y + R_H + GAP
    C_H = CONTENT_H * 0.280    # ≈ 10.9"
    add_rect(slide, cx, C_Y, col_w, C_H, GRAY_LT, TEAL, 1.5)
    section_header(slide, cx, C_Y, col_w, SH, "מסקנות", TEAL)

    c_card_h = (C_H - SH - 0.25) / 4 - 0.12
    for ci, (icon, title, body, bg, accent) in enumerate([
        ("🏆", "Genetic Programming",
         "השיטה הטובה ביותר — 1.59× ב-RA עם 4 features. "
         "מגלה אינטראקציות (mono_pct/neut_pct) שה-LLM לא מצא.",
         ORANGE_LT, ORANGE),
        ("⚡", "LLM (Med-Gemma 4B)",
         "נוסחה: (hct−hgb)/(plt+0.01) — 3 features, 2 ops, lift 1.45×. "
         "ידע רפואי מייצר נוסחאות פשוטות ויעילות.",
         TEAL_LIGHT, TEAL),
        ("📊", "הפריימוורק עובד",
         "השוואה הוגנת: אותם נתונים, split ומטריקות לכל 4 שיטות. "
         "כל שיטה חושפת היבט שונה של המרחב.",
         GRAY_LT, NAVY),
        ("⚠", "מגבלות",
         "CI רחב בשל test set קטן (~300 cases). "
         "הפרשים בין שיטות אינם מובהקים סטטיסטית.",
         GRAY_LT, GRAY),
    ]):
        cy_ = C_Y + SH + 0.15 + ci * (c_card_h + 0.12)
        card(slide, cx + 0.20, cy_, col_w - 0.40, c_card_h,
             icon, title, body, bg=bg, accent=accent,
             title_size=22, body_size=18)

    # ── Section: ולידציה חיצונית + המשך מחקר ────────────────────────────
    V_Y = C_Y + C_H + GAP
    V_H = CONTENT_H - (V_Y - Y0)
    add_rect(slide, cx, V_Y, col_w, V_H, GRAY_LT, NAVY, 1.5)
    section_header(slide, cx, V_Y, col_w, SH,
                   "ולידציה חיצונית  ·  המשך מחקר", NAVY)

    # External validation
    ext_y = V_Y + SH + 0.18
    add_rect(slide, cx + 0.20, ext_y, col_w - 0.40, 0.50, TEAL)
    add_textbox(slide, cx + 0.28, ext_y + 0.07, col_w - 0.56, 0.38,
                "ולידציה חיצונית: EHRSHOT (Stanford) + NHANES",
                21, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    ext_item_h = 1.10
    for ei, (text, dot) in enumerate([
        ("GP מחזיק ביצועים ב-EHRSHOT: RA 1.59× → 2.46× (שיפור!)",   ORANGE),
        ("LLM overfits: Lupus 3.70× (MIMIC) → 1.76× (EHRSHOT)",      TEAL),
        ("NHANES: AUC-PR גולמי ל-RA + Psoriasis (Lift חסר)",          GRAY),
    ]):
        ey = ext_y + 0.62 + ei * ext_item_h
        add_rect(slide, cx + 0.20, ey + 0.15, 0.12, 0.65, dot)
        add_textbox(slide, cx + 0.42, ey + 0.05, col_w - 0.62, ext_item_h - 0.08,
                    text, 18, color=NAVY, align=PP_ALIGN.RIGHT)

    # Future work
    fw_y = ext_y + 0.62 + 3 * ext_item_h + 0.20
    add_rect(slide, cx + 0.20, fw_y, col_w - 0.40, 0.50, NAVY)
    add_textbox(slide, cx + 0.28, fw_y + 0.06, col_w - 0.56, 0.40,
                "🔭  המשך מחקר", 21, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    future = [
        ("Nested Cross-Validation",  "5-fold stratified CV — הערכה מהימנה יותר"),
        ("Seeded GP מכל המחלות",     "LLM seeds כאתחול לאוכלוסיית GP (hybrid)"),
        ("Cross-method equivalence", "Pearson r בין score vectors של שיטות שונות"),
        ("הרחבה למחלות נוספות",      "הפריימוורק agnostic — ICD + CBC = מספיק"),
    ]
    fw_item_h = (V_H - SH - 0.18 - 0.50 - 0.62 - 3 * ext_item_h - 0.20 - 0.62) / 4 - 0.10
    for fi, (title, detail) in enumerate(future):
        fy2 = fw_y + 0.62 + fi * (fw_item_h + 0.10)
        add_rect(slide, cx + 0.20, fy2, col_w - 0.40, fw_item_h, GRAY_LT, NAVY, 0.8)
        add_textbox(slide, cx + 0.30, fy2 + 0.06, col_w - 0.60, 0.44,
                    title, 19, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        add_textbox(slide, cx + 0.30, fy2 + 0.46, col_w - 0.60, fw_item_h - 0.50,
                    detail, 17, color=GRAY, align=PP_ALIGN.RIGHT)

    # ── Save ─────────────────────────────────────────────────────────────
    out = "results/poster/פוסטר.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"  Poster {W:.1f}\" x {H:.1f}\"  |  col_w={col_w:.2f}\"  |  content_h={CONTENT_H:.2f}\"")
    print(f"  Images: fw_h={fw_h:.2f}\"  cons_h={cons_h:.2f}\"  ra_h={ra_h:.2f}\"")


if __name__ == "__main__":
    build()
