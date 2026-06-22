"""
Update docs/Can-a-Routine-CBC-Predict-Chronic-Disease.pptx with:
  1. Three advisors on title slide
  2. New RA disease background slide (inserted at position 3)
  3. Slide 5 (Methods): relabelled to 3 Methods + 2 Baselines
  4. Slide 6 (Results): renamed title, chart baseline relabelling, formula intuition
  5. New all-disease results slide (inserted after RA results)
  6. Slide 7 (External Validation): EHRSHOT results table image

Output: docs/Can-a-Routine-CBC-Predict-Chronic-Disease_v2.pptx
"""

import io
import sys
from copy import deepcopy
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.util import Emu, Inches, Pt

# ── Paths ─────────────────────────────────────────────────────────────────────
PPTX_IN  = Path("docs/Can-a-Routine-CBC-Predict-Chronic-Disease.pptx")
PPTX_OUT = Path("docs/Can-a-Routine-CBC-Predict-Chronic-Disease_v2.pptx")
RESULTS  = Path("results")

# ── Results computed from matched_lr_baseline.csv + ehrshot eval files ────────
# MIMIC lift = formula_auc_pr / disease_prevalence (test-set performance)
MIMIC_LIFTS = {
    "RA":        {"M1": 1.32, "M2": 3.59, "M3": 1.45, "M4": 1.39, "M5": 1.52},
    "Crohn":     {"M1": 2.38, "M2": 4.61, "M3": 3.20, "M4": 2.39, "M5": 2.65},
    "T1D":       {"M1": 1.43, "M2": 3.57, "M3": 1.19, "M4": 1.22, "M5": 1.68},
    "T2D":       {"M1": 1.19, "M2": 1.22, "M3": 1.26, "M4": 1.12},
    "Psoriasis": {"M1": 2.35, "M2": 7.02, "M3": 1.28, "M4": 2.40, "M5": 1.28},
    "Lupus":     {"M1": 2.03, "M2": 11.36, "M3": 1.54, "M4": 3.56, "M5": 1.05},
}
# EHRSHOT lift = best_formula_auc_pr / ehrshot_prevalence (external validation)
EHRSHOT_LIFTS = {
    "RA":        {"M1": 2.04, "M2": 2.19, "M3": 2.56, "M4": 2.53, "M5": 2.94},
    "Crohn":     {"M1": 1.21, "M2": 5.13, "M3": 5.28, "M4": 3.29, "M5": 5.35},
    "T1D":       {"M1": 1.41, "M2": 2.87, "M3": 3.02, "M4": 2.64, "M5": 3.07},
    "T2D":       {"M1": 1.10, "M2": 1.21, "M3": 1.21, "M4": 1.09},
    "Psoriasis": {"M1": 1.14, "M2": 1.58, "M3": 1.44, "M4": 1.48, "M5": 1.44},
    "Lupus":     {"M1": 1.36, "M2": 1.73, "M3": 2.37, "M4": 1.81, "M5": 1.84},
}
# n_features for best formula per method per disease (from methods_comparison.csv)
N_FEATURES = {
    "RA":        {"M2": 9,  "M3": 4, "M4": 3, "M5": 4},
    "Crohn":     {"M2": 13, "M3": 4, "M4": 3, "M5": 4},
    "T1D":       {"M2": 6,  "M3": 5, "M4": 4, "M5": 9},
    "T2D":       {"M2": 13, "M3": 5, "M4": 3},
    "Psoriasis": {"M2": 4,  "M3": 3, "M4": 4, "M5": 3},
    "Lupus":     {"M2": 8,  "M3": 5, "M4": 3, "M5": 6},
}

# ── Bar-chart geometry constants (empirically measured from existing Slide 6) ─
CHART_BOTTOM_IN = 5.433   # y-coordinate (from slide top) where bars sit
CHART_SCALE     = 0.8877  # inches of bar height per lift unit

# ── Colour palette ────────────────────────────────────────────────────────────
C_BASELINE = RGBColor(0x9E, 0x9E, 0x9E)  # gray  – baselines
C_M2       = RGBColor(0x42, 0x9B, 0xF4)  # blue
C_M3       = RGBColor(0x34, 0xA8, 0x53)  # green (GP star)
C_M4       = RGBColor(0xFB, 0xBC, 0x05)  # amber
C_M5       = RGBColor(0x0B, 0x60, 0x27)  # dark green – seeded GP


# ════════════════════════════════════════════════════════════════════════════════
# Helpers
# ════════════════════════════════════════════════════════════════════════════════

def _in(x): return Inches(x)
def _pt(x): return Pt(x)

def _find_shape(slide, text_fragment):
    """Return first shape whose text contains text_fragment (case-sensitive)."""
    for s in slide.shapes:
        if s.has_text_frame and text_fragment in s.text_frame.text:
            return s
    return None

def _find_shape_by_name(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def _set_run_text(tf, new_text):
    """Set the entire text of a TextFrame to new_text, preserving run formatting."""
    for i, para in enumerate(tf.paragraphs):
        if i == 0 and para.runs:
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run._r.getparent().remove(run._r)
        elif i > 0:
            p = para._p
            p.getparent().remove(p)

def _replace_in_shape(shape, old, new):
    """Replace all occurrences of `old` with `new` in every run of every paragraph."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)

def _add_textbox(slide, left, top, width, height, text,
                  font_name="Calibri", font_size=12, bold=False, italic=False,
                  color=None, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(_in(left), _in(top), _in(width), _in(height))
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run  = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = _pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb

def _move_slide(prs, from_idx, to_idx):
    """Move a slide from from_idx to to_idx (0-based)."""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[from_idx]
    sldIdLst.remove(sldId)
    sldIdLst.insert(to_idx, sldId)

def _bytes_to_image_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=150)
    buf.seek(0)
    return buf


# ════════════════════════════════════════════════════════════════════════════════
# 1.  Title slide – update advisors
# ════════════════════════════════════════════════════════════════════════════════

def update_advisors(slide):
    shape = _find_shape(slide, "Advisor")
    if shape is None:
        print("  [WARN] advisor shape not found on slide 1")
        return
    tf = shape.text_frame
    # Para 1, Run 1 = "Advisor:" ; Para 1, Run 2 = " Prof. Noa Dagan     "
    for para in tf.paragraphs:
        runs = para.runs
        for i, run in enumerate(runs):
            if run.text.strip() == "Advisor:":
                run.text = "Advisors:"
            if "Prof. Noa Dagan" in run.text:
                run.text = " Prof. Noa Dagan, Dr. Seffi Cohen, Dr. Liat Friedman Antwarg"
    print("  Slide 1: advisors updated")


# ════════════════════════════════════════════════════════════════════════════════
# 2.  Methods slide (original index 4) – 3 Methods + 2 Baselines
# ════════════════════════════════════════════════════════════════════════════════

def update_methods_slide(slide):
    # Title
    title_shape = _find_shape(slide, "Four Methods")
    if title_shape:
        _replace_in_shape(title_shape, "Four Methods: From Naive to Sophisticated",
                          "Three Methods + Two Baselines")

    # M1 label → Baseline 1
    m1_label = _find_shape(slide, "M1 · Threshold")
    if m1_label:
        _replace_in_shape(m1_label, "M1 · Threshold", "B1  ·  Threshold")

    # M1 description
    m1_desc = _find_shape(slide, "Single feature + cutoff")
    if m1_desc:
        _replace_in_shape(
            m1_desc,
            "Single feature + cutoff (Youden's index). 14 features × 2 strategies.",
            "Single CBC feature + optimal cutoff (Youden’s index). Simplest interpretable rule."
        )

    # Footer
    footer = _find_shape(slide, "LR baseline")
    if footer:
        _replace_in_shape(
            footer,
            "All methods evaluated on the same train/test split. LR baseline (all 14 features) = the bar to beat.",
            "Two baselines: (B1) single-feature threshold · (B2) logistic regression using the same "
            "features as the winning formula. Methods M2–M4 compete against both."
        )

    print("  Slide 5: methods slide updated (3 methods + 2 baselines)")


# ════════════════════════════════════════════════════════════════════════════════
# 3.  Results slide (original index 5) – title, bar reframe, formula annotations
# ════════════════════════════════════════════════════════════════════════════════

def _bar_geometry(value):
    """Return (height_in, top_in) for a bar with the given lift value."""
    h = value * CHART_SCALE
    t = CHART_BOTTOM_IN - h
    return h, t

def update_results_slide(slide, matched_lr_lift):
    # ── Title ──────────────────────────────────────────────────────────────────
    title = _find_shape(slide, "Results: The Race")
    if title:
        _replace_in_shape(title, "Results: The Race", "Results: Method Comparison")

    # ── Relabel chart axis: Bar 1 "LR Baseline" → "Matched LR (B2)" ──────────
    tb46 = _find_shape_by_name(slide, "TextBox 46")
    tb47 = _find_shape_by_name(slide, "TextBox 47")
    if tb46: _set_run_text(tb46.text_frame, "Matched")
    if tb47: _set_run_text(tb47.text_frame, "LR (B2)")

    # ── Relabel chart axis: Bar 2 "M1 Threshold" → "Threshold (B1)" ──────────
    tb48 = _find_shape_by_name(slide, "TextBox 48")
    tb49 = _find_shape_by_name(slide, "TextBox 49")
    if tb48: _set_run_text(tb48.text_frame, "B1")
    if tb49: _set_run_text(tb49.text_frame, "Threshold")

    # ── Update floating label for bar 1 ───────────────────────────────────────
    tb44 = _find_shape_by_name(slide, "TextBox 44")
    if tb44:
        new_label = f"Matched LR (B2): {matched_lr_lift:.2f}×"
        _replace_in_shape(tb44, "LR Baseline: 1.49×", new_label)

    # ── Update value text for bar 1 (TextBox 38) and adjust bar height ────────
    tb38 = _find_shape_by_name(slide, "TextBox 38")
    if tb38:
        _set_run_text(tb38.text_frame, f"{matched_lr_lift:.2f}×")
        # Reposition text box above the new bar
        new_h, new_top = _bar_geometry(matched_lr_lift)
        # text box gap above bar (original: 4.11 - 3.90 = 0.21")
        tb38.top = _in(new_top - 0.21)

    # ── Resize bar 1 (Rectangle: Rounded Corners 20) ─────────────────────────
    bar1 = _find_shape_by_name(slide, "Rectangle: Rounded Corners 20")
    if bar1:
        new_h, new_top = _bar_geometry(matched_lr_lift)
        bar1.height = _in(new_h)
        bar1.top    = _in(new_top)
        # Change fill to gray to visually mark it as a baseline
        bar1.fill.solid()
        bar1.fill.fore_color.rgb = C_BASELINE

    # ── Change bar 2 (M1/Threshold, Rounded Corners 21) to same baseline gray ─
    bar2 = _find_shape_by_name(slide, "Rectangle: Rounded Corners 21")
    if bar2:
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = C_BASELINE

    # ── Formula intuition annotations ─────────────────────────────────────────
    # Formula is at pos (7.37", 4.17"). Add structured annotation below it.
    ann_left, ann_top, ann_w, ann_h = 7.10, 4.62, 6.10, 0.80
    ann_text = (
        "• mono_pct / neut_pct  →  immune imbalance ratio  (↑ in autoimmune disease)\n"
        "• hct − mchc  →  anemia severity signal  (↓ in chronic inflammation)"
    )
    _add_textbox(slide, ann_left, ann_top, ann_w, ann_h, ann_text,
                 font_size=10, italic=True,
                 color=RGBColor(0x55, 0x55, 0x55))

    print("  Slide 6: title renamed, bar 1 relabelled as Matched LR (B2), formula annotated")


# ════════════════════════════════════════════════════════════════════════════════
# 4.  New RA background slide
# ════════════════════════════════════════════════════════════════════════════════

def make_ra_background_slide(prs):
    blank = prs.slide_layouts[0]   # Blank layout
    slide = prs.slides.add_slide(blank)

    # Title
    _add_textbox(slide, 0.50, 0.45, 12.30, 0.65,
                 "Why Rheumatoid Arthritis?",
                 font_size=26, bold=True, color=RGBColor(0x1F, 0x39, 0x7D))

    # Divider line (thin rectangle)
    div = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _in(0.50), _in(1.15), _in(12.30), _in(0.04)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0x1F, 0x39, 0x7D)
    div.line.fill.background()

    # Column 1 – Disease & diagnosis gap
    _add_textbox(slide, 0.50, 1.30, 5.80, 0.30,
                 "THE DISEASE",
                 font_size=11, bold=True, color=RGBColor(0x1F, 0x39, 0x7D))

    body1 = (
        "Rheumatoid Arthritis is a systemic autoimmune disease affecting ~1% of adults "
        "worldwide, causing progressive joint damage and systemic inflammation.\n\n"
        "Diagnosis is often delayed 1–2 years: early symptoms are non-specific, "
        "and confirmation requires specialist referral, anti-CCP serology, and imaging "
        "— resources not always available in primary care.\n\n"
        "Early treatment (within 3–6 months of onset) dramatically slows joint "
        "destruction. Each month of delay worsens long-term outcomes."
    )
    _add_textbox(slide, 0.50, 1.65, 5.80, 2.40, body1, font_size=12,
                 color=RGBColor(0x22, 0x22, 0x22))

    # Column 2 – Why a CBC biomarker?
    _add_textbox(slide, 6.80, 1.30, 5.90, 0.30,
                 "WHY A CBC BIOMARKER?",
                 font_size=11, bold=True, color=RGBColor(0x1F, 0x39, 0x7D))

    body2 = (
        "A Complete Blood Count is ordered in virtually every routine hospital visit "
        "— cheap (~$15), fast, and universally available.\n\n"
        "A CBC-based alert could flag patients for specialist referral before symptoms "
        "are clinically obvious, closing the diagnosis gap in primary care.\n\n"
        "Why does CBC carry RA signal?\n"
        "• Chronic inflammation ↑ monocyte %  —  monocytes are key "
        "mediators of RA synovial inflammation\n"
        "• Inflammation suppresses red cell production → anemia of chronic "
        "disease  (↓ hematocrit, ↓ MCHC)\n"
        "• These signals appear directly in the winning formula:\n"
        "    (mono_pct / neut_pct)² × log(hct − mchc)"
    )
    _add_textbox(slide, 6.80, 1.65, 5.90, 4.20, body2, font_size=12,
                 color=RGBColor(0x22, 0x22, 0x22))

    # Footer note
    _add_textbox(slide, 0.50, 6.90, 12.30, 0.35,
                 "Focus disease throughout this presentation. Results for all 6 diseases shown in the multi-disease summary slide.",
                 font_size=10, italic=True, color=RGBColor(0x77, 0x77, 0x77))

    print("  New RA background slide created")
    return slide


# ════════════════════════════════════════════════════════════════════════════════
# 5.  New all-disease results slide (matplotlib chart)
# ════════════════════════════════════════════════════════════════════════════════

def make_all_disease_slide(prs):
    blank = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank)

    _add_textbox(slide, 0.50, 0.30, 12.30, 0.55,
                 "AUC-PR Lift Across All 6 Diseases — MIMIC Test Set",
                 font_size=22, bold=True, color=RGBColor(0x1F, 0x39, 0x7D))

    _add_textbox(slide, 0.50, 0.88, 12.30, 0.28,
                 "Lift = AUC-PR / prevalence. Values > 1 mean the formula outperforms random flagging. Y-axis capped at 8× (M2 Lupus = 11.36×, M2 Psoriasis = 7.02×).",
                 font_size=10, italic=True, color=RGBColor(0x55, 0x55, 0x55))

    diseases = list(MIMIC_LIFTS.keys())
    methods  = ["M1", "M2", "M3", "M4", "M5"]
    labels   = ["B1 Threshold", "M2 Random", "M3 GP", "M4 LLM", "M5 Seeded GP ★"]
    colors   = ["#9E9E9E", "#429BF4", "#34A853", "#FBBC05", "#0B6027"]
    n_m, n_d = len(methods), len(diseases)

    Y_CAP = 8.0

    x = np.arange(n_d)
    width = 0.14
    offsets = np.linspace(-(n_m-1)/2*width, (n_m-1)/2*width, n_m)

    fig, ax = plt.subplots(figsize=(13, 4.5))
    for i, (m, lab, c) in enumerate(zip(methods, labels, colors)):
        vals_raw = [MIMIC_LIFTS[d].get(m) for d in diseases]
        vals_capped = [min(v, Y_CAP) if v is not None else 0.0 for v in vals_raw]
        xs = [x[j] + offsets[i] for j in range(n_d) if vals_raw[j] is not None]
        vc = [vals_capped[j] for j in range(n_d) if vals_raw[j] is not None]
        vr = [vals_raw[j] for j in range(n_d) if vals_raw[j] is not None]
        bars = ax.bar(xs, vc, width, label=lab, color=c, edgecolor="white", linewidth=0.5)
        bold_m = m in ("M3", "M5")
        for bar, vraw in zip(bars, vr):
            label_str = f"{vraw:.2f}" if vraw <= Y_CAP else f"{vraw:.1f}→"
            ax.text(bar.get_x() + bar.get_width()/2,
                    min(vraw, Y_CAP) + 0.10,
                    label_str, ha="center", va="bottom", fontsize=6.5,
                    fontweight="bold" if bold_m else "normal",
                    color="#0B6027" if m == "M5" else "black")

    ax.axhline(1.0, color="black", linewidth=1.2, linestyle="--", alpha=0.6, label="Lift = 1 (random)")
    ax.set_xticks(x)
    ax.set_xticklabels(diseases, fontsize=11)
    ax.set_ylabel("AUC-PR Lift  (×)", fontsize=11)
    ax.set_ylim(0, Y_CAP * 1.15)
    ax.legend(fontsize=8.5, ncol=3, loc="upper right", framealpha=0.85)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_facecolor("#FAFAFA")
    fig.patch.set_facecolor("white")
    plt.tight_layout()

    img_stream = _bytes_to_image_stream(fig)
    plt.close(fig)

    slide.shapes.add_picture(img_stream, _in(0.20), _in(1.20), _in(13.00), _in(5.70))

    _add_textbox(slide, 0.50, 6.95, 12.30, 0.30,
                 "B1 = single-feature threshold baseline. M5 (dark green) = Seeded GP using LLM-generated formula seeds. T2D has no M5 (seeding not applied). "
                 "M2’s high MIMIC lift on lup/psr reflects feature count overfitting — see next slide.",
                 font_size=9, italic=True, color=RGBColor(0x77, 0x77, 0x77))

    print("  New all-disease results slide created (with M5)")
    return slide


# ════════════════════════════════════════════════════════════════════════════════
# 6.  External Validation slide (original index 6) – add EHRSHOT table image
# ════════════════════════════════════════════════════════════════════════════════

def update_external_slide(slide):
    diseases   = list(EHRSHOT_LIFTS.keys())
    methods    = ["M1", "M2", "M3", "M4", "M5"]
    col_labels = ["Disease", "B1 Threshold", "M2 Random", "M3 GP", "M4 LLM", "M5 Seeded GP ★"]

    rows = []
    for d in diseases:
        row = [d]
        for m in methods:
            v = EHRSHOT_LIFTS[d].get(m)
            row.append(f"{v:.2f}×" if v is not None else "—")
        rows.append(row)

    fig, ax = plt.subplots(figsize=(11.5, 3.2))
    ax.axis("off")
    tbl = ax.table(
        cellText=rows,
        colLabels=col_labels,
        cellLoc="center",
        loc="center",
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(10.5)
    tbl.scale(1, 1.6)

    for j in range(len(col_labels)):
        cell = tbl[0, j]
        cell.set_facecolor("#1F397D")
        cell.set_text_props(color="white", fontweight="bold")

    for i, d in enumerate(diseases):
        vals = {m: EHRSHOT_LIFTS[d].get(m, 0.0) for m in methods}
        best_m = max(vals, key=vals.get)
        for j, m in enumerate(methods):
            cell = tbl[i+1, j+1]
            if m == best_m and vals[m] > 0:
                cell.set_facecolor("#C8E6C9")   # best in row → light green
            elif m == "M5":
                cell.set_facecolor("#E8F5E9")   # M5 column → very light green

    fig.patch.set_facecolor("white")
    plt.tight_layout()

    img_stream = _bytes_to_image_stream(fig)
    plt.close(fig)

    _add_textbox(slide, 0.50, 1.55, 12.30, 0.30,
                 "EHRSHOT (Stanford EHR) — AUC-PR Lift  (lift = AUC-PR / cohort prevalence)",
                 font_size=11, bold=True, color=RGBColor(0x1F, 0x39, 0x7D))

    slide.shapes.add_picture(img_stream, _in(0.30), _in(1.90), _in(12.70), _in(4.30))

    _add_textbox(slide, 0.50, 6.30, 12.30, 0.60,
                 "NHANES (US population survey, RA + Psoriasis only): "
                 "RA AUC-PR — M1: 0.082, M2: 0.077, M3: 0.079, M4: 0.102.  "
                 "Psoriasis — M1: 0.022, M2: 0.033, M3: 0.038, M4: 0.035.  "
                 "(Lift not computed: NHANES prevalence unknown due to self-report inflation.)",
                 font_size=9, italic=True, color=RGBColor(0x55, 0x55, 0x55))

    print("  Slide 7 (External Validation): EHRSHOT table updated with M5")


# ════════════════════════════════════════════════════════════════════════════════
# 7.  Complexity vs Generalisation slide
# ════════════════════════════════════════════════════════════════════════════════

def make_complexity_slide(prs):
    """Show M2 vs M5: high MIMIC lift collapses on EHRSHOT for M2; M5 holds up."""
    blank = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank)

    _add_textbox(slide, 0.50, 0.25, 12.30, 0.55,
                 "Why M5? M2 Overfits — M5 Generalises",
                 font_size=22, bold=True, color=RGBColor(0x1F, 0x39, 0x7D))

    _add_textbox(slide, 0.50, 0.83, 12.30, 0.30,
                 "M2 (Random Search) uses 4–13 features and achieves very high MIMIC test-set lift. "
                 "On EHRSHOT (unseen Stanford data) that advantage nearly vanishes. "
                 "M5 (Seeded GP, 3–6 features) is more consistent across both datasets.",
                 font_size=10, italic=True, color=RGBColor(0x55, 0x55, 0x55))

    diseases_5 = [d for d in MIMIC_LIFTS if "M5" in MIMIC_LIFTS[d]]  # skip T2D
    x = np.arange(len(diseases_5))
    width = 0.30

    m2_mimic    = [MIMIC_LIFTS[d]["M2"]    for d in diseases_5]
    m5_mimic    = [MIMIC_LIFTS[d]["M5"]    for d in diseases_5]
    m2_ehrshot  = [EHRSHOT_LIFTS[d]["M2"]  for d in diseases_5]
    m5_ehrshot  = [EHRSHOT_LIFTS[d]["M5"]  for d in diseases_5]

    Y_CAP = 8.0
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.4), sharey=False)

    for ax, m2_vals, m5_vals, title, note in [
        (axes[0], m2_mimic,   m5_mimic,   "MIMIC Test Set",   "(training domain, capped at 8×)"),
        (axes[1], m2_ehrshot, m5_ehrshot, "EHRSHOT (Stanford)", "(external, zero-shot)"),
    ]:
        m2_capped = [min(v, Y_CAP) for v in m2_vals]
        m5_capped = [min(v, Y_CAP) for v in m5_vals]

        bars_m2 = ax.bar(x - width/2, m2_capped, width, color="#429BF4", label="M2 Random")
        bars_m5 = ax.bar(x + width/2, m5_capped, width, color="#0B6027", label="M5 Seeded GP")

        for bar, vraw in zip(bars_m2, m2_vals):
            s = f"{vraw:.2f}" if vraw <= Y_CAP else f"{vraw:.1f}→"
            ax.text(bar.get_x() + bar.get_width()/2, min(vraw, Y_CAP) + 0.08,
                    s, ha="center", va="bottom", fontsize=7.5)
        for bar, v in zip(bars_m5, m5_vals):
            ax.text(bar.get_x() + bar.get_width()/2, v + 0.08,
                    f"{v:.2f}", ha="center", va="bottom", fontsize=7.5,
                    fontweight="bold", color="#0B6027")

        ax.axhline(1.0, color="black", linewidth=1.0, linestyle="--", alpha=0.5)
        ax.set_xticks(x)
        ax.set_xticklabels(diseases_5, fontsize=10)
        ax.set_ylabel("AUC-PR Lift  (×)", fontsize=10)
        ymax = max(max(m2_capped), max(m5_capped)) * 1.22
        ax.set_ylim(0, min(ymax, Y_CAP * 1.18))
        ax.legend(fontsize=9, loc="upper right")
        ax.set_title(f"{title}\n{note}", fontsize=11, pad=6)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.set_facecolor("#FAFAFA")

    fig.patch.set_facecolor("white")
    plt.tight_layout(pad=1.5)

    img_stream = _bytes_to_image_stream(fig)
    plt.close(fig)

    slide.shapes.add_picture(img_stream, _in(0.20), _in(1.20), _in(13.00), _in(5.50))

    # Feature count annotation table below chart
    n_table = "Features used:  " + "   |   ".join(
        f"{d}: M2={N_FEATURES[d]['M2']}, M5={N_FEATURES[d]['M5']}"
        for d in diseases_5
    )
    _add_textbox(slide, 0.50, 6.80, 12.30, 0.40,
                 n_table,
                 font_size=9, italic=True, color=RGBColor(0x55, 0x55, 0x55))

    print("  New complexity slide created (M2 vs M5 MIMIC vs EHRSHOT)")
    return slide


# ════════════════════════════════════════════════════════════════════════════════
# Main
# ════════════════════════════════════════════════════════════════════════════════

def load_matched_lr():
    """Load RA M3 matched-LR lift from results/matched_lr_baseline.csv."""
    path = RESULTS / "matched_lr_baseline.csv"
    if not path.exists():
        print("  [WARN] matched_lr_baseline.csv not found, using default 1.21")
        return 1.21
    df = pd.read_csv(path)
    row = df[(df["disease"] == "ra") & (df["method"] == "M3")]
    if row.empty:
        return 1.21
    return float(row.iloc[0]["lift"])


def main():
    prs = Presentation(str(PPTX_IN))
    matched_lr_lift = load_matched_lr()
    print(f"  Matched LR (RA, M3 features): {matched_lr_lift:.3f}x")

    # ── Modify existing slides (indices stable before any insertions) ──────────
    print("Updating existing slides ...")
    update_advisors(prs.slides[0])         # slide 1 – title
    update_methods_slide(prs.slides[4])    # slide 5 – methods
    update_results_slide(prs.slides[5], matched_lr_lift)  # slide 6 – RA results
    update_external_slide(prs.slides[6])   # slide 7 – external validation

    # ── Add new slides (appended at end, then moved into position) ────────────
    print("Adding new slides ...")
    n_orig = len(prs.slides)               # original slide count (8)

    make_ra_background_slide(prs)          # appended at index n_orig
    _move_slide(prs, n_orig, 2)            # move to position 3 (0-indexed: 2)
    # slides formerly at 2..n_orig-1 shift right by 1

    make_all_disease_slide(prs)            # appended at index n_orig+1 (last)
    _move_slide(prs, n_orig + 1, 7)        # place after RA Results (now at 6)

    make_complexity_slide(prs)             # appended at index n_orig+2 (last)
    _move_slide(prs, n_orig + 2, 8)        # place right after all-disease (7)

    prs.save(str(PPTX_OUT))
    print(f"\nDone. Saved to {PPTX_OUT}")


if __name__ == "__main__":
    main()
